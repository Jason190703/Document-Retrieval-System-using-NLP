import streamlit as st
import zipfile
import os
import tempfile
from io import BytesIO

from docx import Document
from PyPDF2 import PdfReader
from pptx import Presentation

def extract_zip(zip_file):
    tmpdir = tempfile.TemporaryDirectory()
    with zipfile.ZipFile(zip_file, 'r') as zip_ref:
        zip_ref.extractall(tmpdir.name)
    return tmpdir

def read_txt(file_path):
    try:
        with open(file_path, 'r', encoding='utf-8') as f:
            return f.read()
    except:
        return ""

def read_docx(file_path):
    try:
        doc = Document(file_path)
        full_text = []
        for para in doc.paragraphs:
            full_text.append(para.text)
        return "\n".join(full_text)
    except:
        return ""

def read_pdf(file_path):
    try:
        reader = PdfReader(file_path)
        full_text = []
        for page in reader.pages:
            full_text.append(page.extract_text() or "")
        return "\n".join(full_text)
    except:
        return ""

def read_pptx(file_path):
    try:
        prs = Presentation(file_path)
        full_text = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text"):
                    full_text.append(shape.text)
        return "\n".join(full_text)
    except:
        return ""

def read_file_contents(file_path):
    ext = os.path.splitext(file_path)[1].lower()
    if ext == ".txt":
        return read_txt(file_path)
    elif ext == ".docx":
        return read_docx(file_path)
    elif ext == ".pdf":
        return read_pdf(file_path)
    elif ext == ".pptx":
        return read_pptx(file_path)
    else:
        return ""

def search_keywords_in_text(text, keywords):
    text_lower = text.lower()
    return any(keyword.lower() in text_lower for keyword in keywords)

def search_in_folder(folder_path, keywords):
    matched_files = []
    for root, dirs, files in os.walk(folder_path):
        for file in files:
            full_path = os.path.join(root, file)
            text = read_file_contents(full_path)
            if text and search_keywords_in_text(text, keywords):
                matched_files.append(full_path)
    return matched_files

def main():
    st.title("NLP-Based Intelligent Document Retriever")

    uploaded_zip = st.file_uploader("Upload ZIP file containing documents", type=["zip"])
    keywords_input = st.text_input("Enter keywords to search (comma separated)")

    if uploaded_zip and keywords_input:
        keywords = [k.strip() for k in keywords_input.split(",") if k.strip()]
        with st.spinner("Extracting and searching files..."):
            tmpdir = extract_zip(uploaded_zip)
            matched_files = search_in_folder(tmpdir.name, keywords)

        if matched_files:
            st.write("### Found matches:")
            for file_path in matched_files:
                filename = os.path.basename(file_path)
                col1, col2 = st.columns([3,1])
                with col1:
                    st.write(filename)
                with col2:
                    with open(file_path, "rb") as f:
                        file_bytes = f.read()
                        st.download_button(
                            label="Download",
                            data=file_bytes,
                            file_name=filename,
                            mime="application/octet-stream",
                        )
        else:
            st.warning("No matching files found.")

if __name__ == "__main__":
    main()
